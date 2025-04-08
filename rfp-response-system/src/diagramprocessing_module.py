"""

import os
import uuid
import json
import logging
from dataclasses import dataclass, asdict
from enum import Enum
from typing import List, Dict, Any, Optional, Tuple, Union
from pathlib import Path

import numpy as np
import cv2
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import requests

# For vector storage and search
import faiss
from sentence_transformers import SentenceTransformer

# Configurable LLM client
from llm_client import LLMClient, TogetherAIClient, PrivateLLMClient

# Logging configuration
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class DiagramType(Enum):
    """Enumeration of supported diagram types."""
    FLOWCHART = 'flowchart'
    ARCHITECTURE = 'architecture'
    ORGANIZATIONAL = 'organizational'
    ENTITY_RELATIONSHIP = 'entity_relationship'
    UML = 'uml'
    NETWORK = 'network'
    TIMELINE = 'timeline'
    CHART = 'chart'  # For graphs, pie charts, etc.
    UNKNOWN = 'unknown'

@dataclass
class DiagramContext:
    """Stores the context information for a diagram."""
    preceding_text: str = ""
    following_text: str = ""
    page_number: int = 0
    caption: str = ""
    title: str = ""
    section_title: str = ""

@dataclass
class DiagramMetadata:
    """Metadata associated with an extracted diagram."""
    diagram_id: str
    document_id: str
    diagram_type: DiagramType
    confidence_score: float
    extracted_text: str = ""
    description: str = ""
    tags: List[str] = None
    context: DiagramContext = None
    dimensions: Tuple[int, int] = None
    creation_timestamp: str = ""
    last_updated_timestamp: str = ""
    
    def __post_init__(self):
        if self.tags is None:
            self.tags = []
        if self.context is None:
            self.context = DiagramContext()

class DiagramStorage:
    """Handles storage and retrieval of diagram files and metadata."""
    
    def __init__(self, base_storage_path: str):
        """
        Initialize the diagram storage system.
        
        Args:
            base_storage_path: Base path for storing diagram files and metadata
        """
        self.base_path = Path(base_storage_path)
        self.images_path = self.base_path / "images"
        self.metadata_path = self.base_path / "metadata"
        self.vector_path = self.base_path / "vectors"
        
        # Create directories if they don't exist
        self.images_path.mkdir(parents=True, exist_ok=True)
        self.metadata_path.mkdir(parents=True, exist_ok=True)
        self.vector_path.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"Initialized DiagramStorage at {self.base_path}")
    
    def save_diagram_image(self, diagram_id: str, image: Union[np.ndarray, Image.Image], 
                          format: str = "PNG") -> str:
        """
        Save a diagram image to storage.
        
        Args:
            diagram_id: Unique identifier for the diagram
            image: The image to store (either as numpy array or PIL Image)
            format: Image format to save as
            
        Returns:
            Path to the saved image
        """
        image_path = self.images_path / f"{diagram_id}.{format.lower()}"
        
        if isinstance(image, np.ndarray):
            # Convert OpenCV image to PIL
            image_pil = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
            image_pil.save(str(image_path), format=format)
        elif isinstance(image, Image.Image):
            image.save(str(image_path), format=format)
        else:
            raise ValueError("Image must be either numpy ndarray or PIL Image")
        
        logger.debug(f"Saved diagram image to {image_path}")
        return str(image_path)
    
    def save_metadata(self, metadata: DiagramMetadata) -> str:
        """
        Save diagram metadata to storage.
        
        Args:
            metadata: DiagramMetadata object to save
            
        Returns:
            Path to the saved metadata file
        """
        metadata_path = self.metadata_path / f"{metadata.diagram_id}.json"
        
        with open(metadata_path, 'w') as f:
            json.dump(asdict(metadata), f, indent=2)
        
        logger.debug(f"Saved diagram metadata to {metadata_path}")
        return str(metadata_path)
    
    def get_diagram_image(self, diagram_id: str) -> Optional[Image.Image]:
        """
        Retrieve a diagram image from storage.
        
        Args:
            diagram_id: Unique identifier for the diagram
            
        Returns:
            PIL Image if found, None otherwise
        """
        # Try different possible extensions
        for ext in ['png', 'jpg', 'jpeg', 'svg', 'pdf']:
            image_path = self.images_path / f"{diagram_id}.{ext}"
            if image_path.exists():
                if ext == 'svg' or ext == 'pdf':
                    # For vector formats, may need special handling
                    logger.warning(f"Vector format {ext} needs special handling")
                    # Placeholder for vector format handling
                    return None
                return Image.open(str(image_path))
        
        logger.warning(f"Diagram image not found for ID {diagram_id}")
        return None
    
    def get_metadata(self, diagram_id: str) -> Optional[DiagramMetadata]:
        """
        Retrieve diagram metadata from storage.
        
        Args:
            diagram_id: Unique identifier for the diagram
            
        Returns:
            DiagramMetadata if found, None otherwise
        """
        metadata_path = self.metadata_path / f"{diagram_id}.json"
        if not metadata_path.exists():
            logger.warning(f"Metadata not found for diagram ID {diagram_id}")
            return None
        
        with open(metadata_path, 'r') as f:
            metadata_dict = json.load(f)
            
        # Convert dictionary back to DiagramMetadata
        metadata = DiagramMetadata(
            diagram_id=metadata_dict['diagram_id'],
            document_id=metadata_dict['document_id'],
            diagram_type=DiagramType(metadata_dict['diagram_type']),
            confidence_score=metadata_dict['confidence_score'],
            extracted_text=metadata_dict['extracted_text'],
            description=metadata_dict['description'],
            tags=metadata_dict['tags'],
            dimensions=tuple(metadata_dict['dimensions']) if metadata_dict.get('dimensions') else None,
            creation_timestamp=metadata_dict['creation_timestamp'],
            last_updated_timestamp=metadata_dict['last_updated_timestamp']
        )
        
        # Reconstruct context
        if context_dict := metadata_dict.get('context'):
            metadata.context = DiagramContext(
                preceding_text=context_dict.get('preceding_text', ''),
                following_text=context_dict.get('following_text', ''),
                page_number=context_dict.get('page_number', 0),
                caption=context_dict.get('caption', ''),
                title=context_dict.get('title', ''),
                section_title=context_dict.get('section_title', '')
            )
            
        return metadata
    
    def list_diagrams(self, document_id: Optional[str] = None) -> List[str]:
        """
        List all diagram IDs, optionally filtered by document ID.
        
        Args:
            document_id: Optional document ID to filter by
            
        Returns:
            List of diagram IDs
        """
        diagram_ids = []
        
        for metadata_file in self.metadata_path.glob("*.json"):
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
                
            if document_id is None or metadata.get('document_id') == document_id:
                diagram_ids.append(metadata.get('diagram_id'))
                
        return diagram_ids

class DiagramSearchEngine:
    """Handles indexing and searching of diagram content and metadata."""
    
    def __init__(self, storage: DiagramStorage, embedding_model: str = "all-MiniLM-L6-v2"):
        """
        Initialize the diagram search engine.
        
        Args:
            storage: DiagramStorage instance to use for retrieving metadata
            embedding_model: Name of the sentence-transformer model to use for embeddings
        """
        self.storage = storage
        self.embedding_dim = 384  # Default for the specified model
        self.embedder = SentenceTransformer(embedding_model)
        
        # Initialize FAISS index
        self.index = faiss.IndexFlatIP(self.embedding_dim)  # Inner product for cosine similarity
        
        # Mapping from index position to diagram ID
        self.index_to_id = []
        
        # Path for saving/loading the index
        self.index_path = storage.vector_path / "diagram_index.faiss"
        self.mapping_path = storage.vector_path / "diagram_id_mapping.json"
        
        # Load existing index if available
        self._load_index()
        
        logger.info(f"Initialized DiagramSearchEngine with model {embedding_model}")
    
    def _load_index(self):
        """Load existing index and ID mapping if available."""
        if self.index_path.exists() and self.mapping_path.exists():
            try:
                self.index = faiss.read_index(str(self.index_path))
                
                with open(self.mapping_path, 'r') as f:
                    self.index_to_id = json.load(f)
                    
                logger.info(f"Loaded existing search index with {len(self.index_to_id)} diagrams")
            except Exception as e:
                logger.error(f"Error loading search index: {e}")
                # Initialize new index
                self.index = faiss.IndexFlatIP(self.embedding_dim)
                self.index_to_id = []
    
    def _save_index(self):
        """Save the current index and ID mapping."""
        try:
            faiss.write_index(self.index, str(self.index_path))
            
            with open(self.mapping_path, 'w') as f:
                json.dump(self.index_to_id, f)
                
            logger.info(f"Saved search index with {len(self.index_to_id)} diagrams")
        except Exception as e:
            logger.error(f"Error saving search index: {e}")
    
    def index_diagram(self, diagram_id: str) -> bool:
        """
        Index a diagram by its ID.
        
        Args:
            diagram_id: ID of the diagram to index
            
        Returns:
            Success status
        """
        metadata = self.storage.get_metadata(diagram_id)
        if not metadata:
            logger.warning(f"Cannot index diagram {diagram_id}: metadata not found")
            return False
        
        # Create a combined text representation for embedding
        text_to_embed = f"{metadata.description} {metadata.extracted_text}"
        if metadata.context:
            text_to_embed += f" {metadata.context.title} {metadata.context.caption}"
        if metadata.tags:
            text_to_embed += f" {' '.join(metadata.tags)}"
            
        # Create embedding
        embedding = self.embedder.encode([text_to_embed])[0]
        embedding = embedding.reshape(1, -1)  # Reshape for FAISS
        
        # Add to index
        self.index.add(np.array(embedding, dtype=np.float32))
        self.index_to_id.append(diagram_id)
        
        # Save updated index
        self._save_index()
        
        logger.info(f"Indexed diagram {diagram_id}")
        return True
    
    def reindex_all(self) -> int:
        """
        Rebuild the entire search index.
        
        Returns:
            Number of diagrams indexed
        """
        # Clear existing index
        self.index = faiss.IndexFlatIP(self.embedding_dim)
        self.index_to_id = []
        
        # Get all diagram IDs
        diagram_ids = self.storage.list_diagrams()
        
        # Index each diagram
        successful = 0
        for diagram_id in diagram_ids:
            if self.index_diagram(diagram_id):
                successful += 1
                
        logger.info(f"Reindexed {successful} out of {len(diagram_ids)} diagrams")
        return successful
    
    def search(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """
        Search for diagrams matching the query.
        
        Args:
            query: Search query
            top_k: Number of results to return
            
        Returns:
            List of result dictionaries with diagram metadata and scores
        """
        # Create query embedding
        query_embedding = self.embedder.encode([query])[0]
        query_embedding = np.array([query_embedding], dtype=np.float32)
        
        # Perform search
        scores, indices = self.index.search(query_embedding, min(top_k, len(self.index_to_id)))
        
        # Build results
        results = []
        for i, idx in enumerate(indices[0]):
            if idx < 0 or idx >= len(self.index_to_id):
                continue
                
            diagram_id = self.index_to_id[idx]
            metadata = self.storage.get_metadata(diagram_id)
            
            if metadata:
                results.append({
                    "diagram_id": diagram_id,
                    "score": float(scores[0][i]),
                    "metadata": asdict(metadata)
                })
        
        return results
    
    def search_by_type(self, diagram_type: DiagramType, top_k: int = 10) -> List[str]:
        """
        Get diagrams of a specific type.
        
        Args:
            diagram_type: Type of diagram to search for
            top_k: Maximum number of results
            
        Returns:
            List of diagram IDs
        """
        matching_ids = []
        
        # This is not using the vector index but doing a metadata scan
        # For production systems, consider maintaining separate indices by type
        for diagram_id in self.storage.list_diagrams():
            metadata = self.storage.get_metadata(diagram_id)
            if metadata and metadata.diagram_type == diagram_type:
                matching_ids.append(diagram_id)
                
            if len(matching_ids) >= top_k:
                break
                
        return matching_ids
    
    def search_by_document(self, document_id: str) -> List[str]:
        """
        Get all diagrams associated with a document.
        
        Args:
            document_id: ID of the document
            
        Returns:
            List of diagram IDs
        """
        return self.storage.list_diagrams(document_id=document_id)

class DiagramExtractor:
    """
    Extracts diagrams from various document formats while preserving context.
    Supports PDF, DOCX, and PPTX formats.
    """
    
    def __init__(self, temp_dir: str = "/tmp/diagram_extraction"):
        """
        Initialize the diagram extractor.
        
        Args:
            temp_dir: Directory for temporary files
        """
        self.temp_dir = Path(temp_dir)
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        logger.info(f"Initialized DiagramExtractor with temp dir {temp_dir}")
        
    def extract_from_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extract diagrams from a PDF file.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of dictionaries containing image data and context
        """
        logger.info(f"Extracting diagrams from PDF: {pdf_path}")
        results = []
        
        try:
            # Open the PDF
            pdf_document = fitz.open(pdf_path)
            
            for page_num, page in enumerate(pdf_document):
                # Extract text for context
                page_text = page.get_text()
                
                # Extract images
                image_list = page.get_images(full=True)
                
                for img_idx, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = pdf_document.extract_image(xref)
                    
                    if not base_image:
                        continue
                        
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # Save image to temporary file
                    temp_img_path = self.temp_dir / f"pdf_img_{page_num}_{img_idx}.{image_ext}"
                    with open(temp_img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Load as PIL image
                    image = Image.open(temp_img_path)
                    
                    # Get image position on page
                    rect = page.get_image_bbox(xref)
                    
                    # Extract context (text before and after the image)
                    # This is a simplification - would need more sophisticated logic for real-world use
                    preceding_text = page_text[:int(len(page_text)/2)]  # Simplified
                    following_text = page_text[int(len(page_text)/2):]  # Simplified
                    
                    # Try to extract caption (text near the image)
                    # This is also simplified - real implementation would need more sophisticated logic
                    caption = self._find_caption_in_text(page_text, rect)
                    
                    results.append({
                        "image": image,
                        "format": image_ext.upper(),
                        "page_number": page_num + 1,
                        "position": {
                            "x1": rect.x0, "y1": rect.y0,
                            "x2": rect.x1, "y2": rect.y1
                        },
                        "context": {
                            "preceding_text": preceding_text,
                            "following_text": following_text,
                            "caption": caption,
                            "page_number": page_num + 1
                        }
                    })
            
            # Clean up
            pdf_document.close()
            
        except Exception as e:
            logger.error(f"Error extracting diagrams from PDF {pdf_path}: {e}")
            
        logger.info(f"Extracted {len(results)} diagrams from PDF {pdf_path}")
        return results
    
    def extract_from_docx(self, docx_path: str) -> List[Dict[str, Any]]:
        """
        Extract diagrams from a DOCX file.
        
        Args:
            docx_path: Path to the DOCX file
            
        Returns:
            List of dictionaries containing image data and context
        """
        logger.info(f"Extracting diagrams from DOCX: {docx_path}")
        results = []
        
        try:
            # Open the document
            doc = docx.Document(docx_path)
            
            image_count = 0
            current_section = ""
            preceding_paragraphs = []
            
            # Process each paragraph to track context
            for para_idx, para in enumerate(doc.paragraphs):
                # Track section headings
                if para.style.name.startswith('Heading'):
                    current_section = para.text
                
                # Add to context buffer
                preceding_paragraphs.append(para.text)
                # Keep only the last 5 paragraphs for context
                if len(preceding_paragraphs) > 5:
                    preceding_paragraphs.pop(0)
                    
            # Extract images from the document
            temp_dir = self.temp_dir / f"docx_{Path(docx_path).stem}"
            temp_dir.mkdir(exist_ok=True)
            
            # Unzip the docx file (it's actually a zip)
            import zipfile
            with zipfile.ZipFile(docx_path) as zip_ref:
                zip_ref.extractall(temp_dir)
            
            # Find and process images
            word_rels_dir = temp_dir / "word" / "media"
            if word_rels_dir.exists():
                for img_path in word_rels_dir.glob("*"):
                    if img_path.suffix.lower() in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
                        image = Image.open(img_path)
                        
                        # Try to find the caption (simplified)
                        caption = self._find_caption_in_docx(doc, image_count)
                        
                        results.append({
                            "image": image,
                            "format": img_path.suffix[1:].upper(),
                            "context": {
                                "preceding_text": "\n".join(preceding_paragraphs),
                                "following_text": "",  # Simplified
                                "caption": caption,
                                "section_title": current_section
                            }
                        })
                        
                        image_count += 1
                        
            # Clean up temp files
            import shutil
            shutil.rmtree(temp_dir)
            
        except Exception as e:
            logger.error(f"Error extracting diagrams from DOCX {docx_path}: {e}")
            
        logger.info(f"Extracted {len(results)} diagrams from DOCX {docx_path}")
        return results
    
    def extract_from_pptx(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        Extract diagrams from a PPTX file.
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            List of dictionaries containing image data and context
        """
        logger.info(f"Extracting diagrams from PPTX: {pptx_path}")
        results = []
        
        try:
            # Open the presentation
            presentation = Presentation(pptx_path)
            
            for slide_idx, slide in enumerate(presentation.slides):
                # Extract slide title
                slide_title = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.shape_type == 1:  # Title
                        slide_title = shape.text
                        break
                
                # Extract slide text for context
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                # Extract images
                temp_dir = self.temp_dir / f"pptx_{Path(pptx_path).stem}_{slide_idx}"
                temp_dir.mkdir(exist_ok=True)
                
                img_count = 0
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        # Get image data
                        image = shape.image
                        image_bytes = image.blob
                        
                        # Save image to temp file
                        img_ext = 'png'  # Default
                        if hasattr(image, 'ext'):
                            img_ext = image.ext
                            
                        temp_img_path = temp_dir / f"slide_{slide_idx}_img_{img_count}.{img_ext}"
                        with open(temp_img_path, "wb") as img_file:
                            img_file.write(image_bytes)
                        
                        # Load as PIL image
                        try:
                            pil_image = Image.open(temp_img_path)
                            
                            results.append({
                                "image": pil_image,
                                "format": img_ext.upper(),
                                "page_number": slide_idx + 1,
                                "context": {
                                    "preceding_text": slide_text,
                                    "following_text": "",
                                    "caption": "",  # PPT doesn't have standard captions
                                    "title": slide_title,
                                    "page_number": slide_idx + 1
                                }
                            })
                            
                            img_count += 1
                        except Exception as e:
                            logger.warning(f"Error loading image: {e}")
                
                # Clean up temp files
                import shutil
                shutil.rmtree(temp_dir)
                
        except Exception as e:
            logger.error(f"Error extracting diagrams from PPTX {pptx_path}: {e}")
            
        logger.info(f"Extracted {len(results)} diagrams from PPTX {pptx_path}")
        return results
    
    def _find_caption_in_text(self, text: str, rect) -> str:
        """Find potential image caption in text (simplified)."""
        # This is a placeholder for more sophisticated logic
        # Real implementation would use position data and text analysis
        
        lines = text.split('\n')
        for line in lines:
            if line.lower().startswith(('figure', 'fig.', 'fig ', 'diagram', 'chart')):
                return line
                
        return ""
    
    def _find_caption_in_docx(self, doc, image_idx: int) -> str:
        """Find potential image caption in Word document (simplified)."""
        # Look for captions (simplified approach)
        for para in doc.paragraphs:
            if para.text.lower().startswith(('figure', 'fig.', 'fig ', 'diagram', 'chart')):
                return para.text
                
        return ""

class DiagramAnalyzer:
    """
    Analyzes extracted diagrams to classify them, extract text, 
    and generate structured representations.
    """
    
    def __init__(self, llm_client: LLMClient, ocr_engine: str = "tesseract"):
        """
        Initialize the diagram analyzer.
        
        Args:
            llm_client: LLM client for diagram analysis
            ocr_engine: OCR engine to use ('tesseract' or 'cloud')
        """
        self.llm_client = llm_client
        self.ocr_engine = ocr_engine
        
        # For diagram classification
        self.classifier = self._initialize_classifier()
        
        logger.info(f"Initialized DiagramAnalyzer with OCR engine {ocr_engine}")
    
    def _initialize_classifier(self):
        """Initialize the diagram classifier model."""
        # In a real implementation, you might use a pre-trained model
        # For this example, we'll use a simplified rule-based approach
        logger.info("Initializing diagram classifier (simplified)")
        return None
    
    def analyze_diagram(self, image: Image.Image, context: Dict[str, Any]) -> Dict[str, Any]:
        """
        Perform comprehensive analysis of a diagram.
        
        Args:
            image: PIL Image of the diagram
            context: Contextual information about the diagram
            
        Returns:
            Dictionary with analysis results
        """
        # Extract text using OCR
        extracted_text = self.perform_ocr(image)
        
        # Classify diagram type
        diagram_type, confidence = self.classify_diagram_type(image, extracted_text, context)
        
        # Generate description using LLM
        description = self.generate_description(image, extracted_text, diagram_type, context)
        
        # Try to generate structured representation
        structured_representation = self.generate_structured_representation(
            image, extracted_text, diagram_type, context
        )
        
        # Generate tags
        tags = self.generate_tags(extracted_text, description, diagram_type, context)
        
        return {
            "extracted_text": extracted_text,
            "diagram_type": diagram_type,
            "confidence_score": confidence,
            "description": description,
            "structured_representation": structured_representation,
            "tags": tags
        }
    
    def perform_ocr(self, image: Image.Image) -> str:
        """
        Extract text from diagram using OCR.
        
        Args:
            image: PIL Image to process
            
        Returns:
            Extracted text string
        """
        if self.ocr_engine == "tesseract":
            # Convert to grayscale for better OCR
            gray_image = image.convert('L')
            
            # Use pytesseract
            try:
                extracted_text = pytesseract.image_to_string(gray_image)
                logger.debug(f"Extracted {len(extracted_text)} characters with OCR")
                return extracted_text
            except Exception as e:
                logger.error(f"OCR error: {e}")
                return ""
        elif self.ocr_engine == "cloud":
            # Placeholder for cloud OCR service integration
            logger.warning("Cloud OCR not implemented, falling back to tesseract")
            return pytesseract.image_to_string(image)
        else:
            logger.error(f"Unknown OCR engine: {self.ocr_engine}")
            return ""
    
    def classify_diagram_type(self, image: Image.Image, text: str, context: Dict[str, Any]) -> Tuple[DiagramType, float]:
        """
        Classify the type of diagram.
        
        Args:
            image: PIL Image of the diagram
            text: Extracted text from OCR
            context: Contextual information
            
        Returns:
            DiagramType enum and confidence score
        """
        # This is a simplified rule-based approach
        # In a real system, you would use a trained classifier
        
        # Look for keywords in the extracted text and context
        text_lower = text.lower()
        context_text = ""
        
        if context.get("caption"):
            context_text += context["caption"].lower() + " "
        if context.get("title"):
            context_text += context["title"].lower() + " "
        if context.get("preceding_text"):
            context_text += context["preceding_text"].lower() + " "
            
        # Simple keyword matching
        if "flow" in text_lower or "flow" in context_text or "process" in text_lower:
            return DiagramType.FLOWCHART, 0.8
        elif "architecture" in text_lower or "architecture" in context_text or "system" in text_lower:
            return DiagramType.ARCHITECTURE, 0.8
        elif "organization" in text_lower or "org chart" in context_text or "reporting" in text_lower:
            return DiagramType.ORGANIZATIONAL, 0.8
        elif "entity" in text_lower or "relationship" in text_lower or "ER diagram" in context_text:
            return DiagramType.ENTITY_RELATIONSHIP, 0.8
        elif "class" in text_lower or "uml" in text_lower or "object" in context_text:
            return DiagramType.UML, 0.8
        elif "network" in text_lower or "topology" in text_lower:
            return DiagramType.NETWORK, 0.7
        elif "timeline" in text_lower or "gantt" in context_text:
            return DiagramType.TIMELINE, 0.7
        elif "chart" in text_lower or "graph" in text_lower or "plot" in context_text:
            return DiagramType.CHART, 0.7
        else:
            # Can't determine type confidently
            return DiagramType.UNKNOWN, 0.5