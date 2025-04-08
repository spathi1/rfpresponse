import re
from typing import Dict, List, Tuple, Optional
import PyPDF2
import docx
import markdown
from pptx import Presentation
import apache_tika
from apache_tika import parser as tika_parser

class DocumentPreprocessor:
    def __init__(self, chunk_strategy='section'):
        """
        Initialize document preprocessor with specified chunking strategy
        
        Args:
            chunk_strategy: Strategy for chunking documents ('section', 'paragraph', 'heading')
        """
        self.chunk_strategy = chunk_strategy
        # Initialize Tika if needed
        apache_tika.initVM()
    
    def preprocess_document(self, document_info: Dict) -> Dict:
        """
        Extract text and structure from document based on its type
        
        Args:
            document_info: Document information dictionary from DocumentInventory
            
        Returns:
            Processed document with extracted text and structure
        """
        doc_path = document_info['path']
        extension = document_info['extension']
        
        # Extract content based on file type
        if extension == '.pdf':
            extracted_content = self._extract_from_pdf(doc_path)
        elif extension == '.docx':
            extracted_content = self._extract_from_docx(doc_path)
        elif extension == '.pptx':
            extracted_content = self._extract_from_pptx(doc_path)
        elif extension == '.md':
            extracted_content = self._extract_from_markdown(doc_path)
        elif extension == '.txt':
            extracted_content = self._extract_from_text(doc_path)
        else:
            # Fallback to Apache Tika for unknown formats
            extracted_content = self._extract_with_tika(doc_path)
        
        # Apply chunking strategy
        chunks = self._chunk_content(extracted_content)
        
        # Update document info with extracted content
        document_info.update({
            'content': extracted_content,
            'chunks': chunks,
            'total_chunks': len(chunks)
        })
        
        return document_info
    
    def _extract_from_pdf(self, path: str) -> Dict:
        """Extract text and structure from PDF document"""
        content = {"full_text": "", "sections": []}
        
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Extract text from each page
            for page_num, page in enumerate(pdf_reader.pages):
                text = page.extract_text()
                content["full_text"] += text
                
                # Try to identify sections based on font size or formatting
                # This is simplified - real implementation would need more sophisticated PDF parsing
                section_matches = re.finditer(r'(?m)^(.*?)(?:\r?\n)', text)
                for match in section_matches:
                    potential_heading = match.group(1).strip()
                    if len(potential_heading) < 100 and potential_heading.isupper():
                        content["sections"].append({
                            "heading": potential_heading,
                            "text": text,
                            "page": page_num + 1
                        })
        
        return content
    
    def _extract_from_docx(self, path: str) -> Dict:
        """Extract text and structure from DOCX document"""
        content = {"full_text": "", "sections": []}
        doc = docx.Document(path)
        
        current_section = {"heading": "", "text": ""}
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            content["full_text"] += text + "\n"
            
            # Check if paragraph is a heading
            if paragraph.style.name.startswith('Heading'):
                # Save previous section if it exists
                if current_section["text"]:
                    content["sections"].append(current_section.copy())
                
                # Start new section
                current_section = {
                    "heading": text,
                    "text": "",
                    "heading_level": int(paragraph.style.name.replace('Heading ', '')) 
                    if paragraph.style.name != 'Heading' else 1
                }
            else:
                # Add to current section
                current_section["text"] += text + "\n"
        
        # Add the last section
        if current_section["text"]:
            content["sections"].append(current_section)
            
        return content
    
    def _extract_from_pptx(self, path: str) -> Dict:
        """Extract text and structure from PPTX document"""
        content = {"full_text": "", "sections": []}
        presentation = Presentation(path)
        
        for slide_num, slide in enumerate(presentation.slides):
            slide_text = ""
            
            # Extract slide title
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if hasattr(shape, 'is_title') and shape.is_title:
                        title = shape.text.strip()
                    
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs)
                        slide_text += paragraph_text + "\n"
            
            content["full_text"] += slide_text
            
            # Add slide as a section
            content["sections"].append({
                "heading": title or f"Slide {slide_num + 1}",
                "text": slide_text,
                "slide_number": slide_num + 1
            })
            
        return content
    
    def _extract_from_markdown(self, path: str) -> Dict:
        """Extract text and structure from Markdown document"""
        content = {"full_text": "", "sections": []}
        
        with open(path, 'r', encoding='utf-8') as file:
            md_text = file.read()
            
        content["full_text"] = md_text
        
        # Use regex to find headings and their content
        heading_pattern = r'(#{1,6})\s+(.+?)(?:\n|$)(.*?)(?=\n#{1,6}\s+|\Z)'
        matches = re.finditer(heading_pattern, md_text, re.DOTALL)
        
        for match in matches:
            heading_level = len(match.group(1))
            heading_text = match.group(2).strip()
            section_text = match.group(3).strip()
            
            content["sections"].append({
                "heading": heading_text,
                "text": section_text,
                "heading_level": heading_level
            })
            
        return content
    
    def _extract_from_text(self, path: str) -> Dict:
        """Extract text from plain text document"""
        with open(path, 'r', encoding='utf-8') as file:
            text = file.read()
            
        content = {
            "full_text": text,
            "sections": [{
                "heading": "",
                "text": text
            }]
        }
        
        return content
    
    def _extract_with_tika(self, path: str) -> Dict:
        """Extract text using Apache Tika for unknown formats"""
        parsed = tika_parser.from_file(path)
        text = parsed["content"]
        
        content = {
            "full_text": text,
            "sections": [{
                "heading": "",
                "text": text
            }]
        }
        
        return content
    
    def _chunk_content(self, content: Dict) -> List[Dict]:
        """Chunk document content based on specified strategy"""
        chunks = []
        
        if self.chunk_strategy == 'section':
            # Chunk by document sections
            for section in content["sections"]:
                chunks.append({
                    "heading": section["heading"],
                    "text": section["text"],
                    "metadata": {k: v for k, v in section.items() 
                               if k not in ["heading", "text"]}
                })
                
        elif self.chunk_strategy == 'paragraph':
            # Chunk by paragraphs
            for section in content["sections"]:
                paragraphs = re.split(r'\n\s*\n', section["text"])
                for i, para in enumerate(paragraphs):
                    if para.strip():
                        chunks.append({
                            "heading": section["heading"],
                            "text": para.strip(),
                            "paragraph_index": i,
                            "metadata": {k: v for k, v in section.items() 
                                       if k not in ["heading", "text"]}
                        })
                        
        elif self.chunk_strategy == 'heading':
            # One chunk per heading and its immediate content
            for section in content["sections"]:
                # Split section by subheadings if they exist
                if "heading_level" in section:
                    current_level = section["heading_level"]
                    subheading_pattern = r'\n(#{1,' + str(current_level+1) + r'})\s+(.+?)(?:\n|$)'
                    parts = re.split(subheading_pattern, section["text"])
                    
                    if len(parts) > 1:
                        # Process parts with subheadings
                        chunks.append({
                            "heading": section["heading"],
                            "text": parts[0].strip(),
                            "metadata": {k: v for k, v in section.items() 
                                       if k not in ["heading", "text"]}
                        })
                        
                        for i in range(1, len(parts), 3):
                            if i+2 < len(parts):
                                sub_level = len(parts[i])
                                sub_heading = parts[i+1]
                                sub_text = parts[i+2]
                                
                                chunks.append({
                                    "heading": sub_heading,
                                    "text": sub_text.strip(),
                                    "parent_heading": section["heading"],
                                    "heading_level": sub_level,
                                    "metadata": {k: v for k, v in section.items() 
                                               if k not in ["heading", "text"]}
                                })
                    else:
                        # No subheadings found
                        chunks.append({
                            "heading": section["heading"],
                            "text": section["text"].strip(),
                            "metadata": {k: v for k, v in section.items() 
                                       if k not in ["heading", "text"]}
                        })
                else:
                    # No heading level info, treat as one chunk
                    chunks.append({
                        "heading": section["heading"],
                        "text": section["text"].strip(),
                        "metadata": {k: v for k, v in section.items() 
                                   if k not in ["heading", "text"]}
                    })
                        
        return chunks